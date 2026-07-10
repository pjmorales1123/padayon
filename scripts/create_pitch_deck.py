#!/usr/bin/env python3
"""Generate PADAYON hackathon pitch deck."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
IMG_DIR = os.path.join(ROOT, "screenshots")
OUT = os.path.join(ROOT, "PADAYON_Pitch_Deck.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    bg.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.75), Inches(2.4), Inches(11.8), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    tb2 = slide.shapes.add_textbox(Inches(0.75), Inches(4.1), Inches(11.8), Inches(1.0))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0x93, 0xC5, 0xFD)
    p2.alignment = PP_ALIGN.CENTER

    tb3 = slide.shapes.add_textbox(Inches(0.75), Inches(6.6), Inches(11.8), Inches(0.5))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "AMD Developer Hackathon: ACT II — Unicorn Track + Best AMD-Hosted Gemma Project"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    p3.alignment = PP_ALIGN.CENTER
    return slide

def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x25, 0x63, 0xEB)
    bg.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(3.0), Inches(11.8), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, bullets, subtitle=None, image_path=None, image_left=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    bg.line.fill.background()

    # Title
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.9))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0x64, 0x70, 0x8B)
        p2.space_before = Pt(6)

    # Content area
    if image_path and os.path.exists(image_path):
        if image_left:
            text_left = Inches(6.3)
            text_width = Inches(6.4)
            img_left = Inches(0.6)
            img_top = Inches(1.6)
            img_w = Inches(5.5)
            img_h = Inches(5.0)
        else:
            text_left = Inches(0.6)
            text_width = Inches(6.4)
            img_left = Inches(6.6)
            img_top = Inches(1.6)
            img_w = Inches(6.0)
            img_h = Inches(5.0)
        slide.shapes.add_picture(image_path, img_left, img_top, width=img_w)
    else:
        text_left = Inches(0.6)
        text_width = Inches(12.1)

    tb = slide.shapes.add_textbox(text_left, Inches(1.5), text_width, Inches(5.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x33, 0x41, 0x55)
        p.space_after = Pt(12)
        p.level = 0
    return slide

# 1. Title
add_title_slide(prs,
    "PADAYON",
    "An AI learning partner that turns messy student notes into organized,\ncurriculum-aligned materials in the language students understand best.")

# 2. The problem
add_content_slide(prs,
    "The Problem: Filipino Students Are Falling Behind",
    [
        "PISA 2022: Philippines ranked 77th out of 81 countries in reading, math, and science.",
        "Scores: Reading 347, Math 355, Science 356 — vs. OECD averages of 472–485.",
        "Academic English is a barrier: Science and Math are taught in English, but many students think best in Cebuano/Filipino.",
        "Students lack organized study habits — notes are scattered, review is last-minute, and study materials are not personalized.",
        "Most AI tools today give instant answers, enabling shortcut learning instead of building real understanding.",
        "Result: learners complete homework but do not master concepts."
    ],
    subtitle="Sources: PISA 2022 Results; Philippine Basic Education; MTB-MLE research on mother-tongue instruction")

# 3. The opportunity
add_content_slide(prs,
    "The Opportunity: AI That Teaches, Not Just Answers",
    [
        "85% of students already use AI for schoolwork — the habit is formed, but the direction is wrong.",
        "Research shows students want AI for explaining concepts and generating ideas, not writing entire outputs.",
        "Mother-tongue-based instruction improves comprehension, confidence, and problem-solving.",
        "What if AI could meet students in their own language, organize their chaos, and grow with them?"
    ],
    subtitle="Stanford GSE / Faculty Focus surveys on AI use; MTB-MLE Philippines studies")

# 4. Solution overview
add_content_slide(prs,
    "PADAYON: Your AI Study Partner",
    [
        "Messy notes, photos, or questions → automatic subject/topic classification.",
        "Aligns every input to the Philippine curriculum / Budget of Work.",
        "Auto-generates clean notes, reviewer, flashcards, quiz, summary, and story.",
        "Teaches through translanguaging: Cebuano or Filipino first, then academic English.",
        "Remembers each learner's strengths, weaknesses, language confidence, and progress.",
        "Retrieves past materials across sessions — learning continues where it left off."
    ],
    image_path=os.path.join(IMG_DIR, "padayon_home.png"),
    image_left=False)

# 5. How it works - agent architecture
add_content_slide(prs,
    "Agent Architecture Built for Learning",
    [
        "Classifier Agent — detects subject, topic, intent, and language from student input.",
        "Curriculum Alignment Agent — matches the topic to DepEd competencies and learning sequence.",
        "Organizer Agent — saves everything into the right subject/topic folder.",
        "Material Creator Agent — builds notes, flashcards, quizzes, reviewers, and stories.",
        "Teaching Agent — explains with translanguaging and guided questions.",
        "Assessment Agent — checks mastery with hints, feedback, and progress scoring.",
        "Memory Agent — updates the learner profile after every interaction."
    ],
    subtitle="Seven real agents orchestrated through Fireworks AI on AMD hardware")

# 6. Translanguaging + memory
add_content_slide(prs,
    "Translanguaging + Adaptive Memory",
    [
        "Students ask in Cebuano: 'Unsa diay ang photosynthesis?'",
        "PADAYON answers in Cebuano first, then introduces the English term naturally.",
        "The learner profile records language confidence, learning style, strengths, and weak areas.",
        "Every chat shapes the next explanation — it feels like a tutor who remembers you.",
        "Mastery score tracks progress per topic: started → developing → mastered."
    ],
    image_path=os.path.join(IMG_DIR, "padayon_profile.png"),
    image_left=True)

# 7. Library
add_content_slide(prs,
    "Automatic Organization in Your Library",
    [
        "Every topic is filed under Subject → Subcategory → Topic.",
        "One-click access to original notes, clean notes, reviewer, flashcards, quiz, summary, and story.",
        "Progress bar and mastery status show what needs review.",
        "Students can return days later and say 'show my flashcards' — PADAYON retrieves instantly."
    ],
    image_path=os.path.join(IMG_DIR, "padayon_library.png"),
    image_left=False)

# 8. Gemma + AMD
add_content_slide(prs,
    "AMD-Ready Architecture with Gemma Support",
    [
        "Production runtime: Fireworks AI API — Fireworks is an AMD partner hosting fast, serverless models on AMD infrastructure.",
        "Default pipeline uses Fireworks serverless models (DeepSeek-V4 / Kimi-K2.5) for reliability and speed.",
        "Demo toggle switches to Gemma 3/4 when an endpoint is available via Fireworks on-demand or AMD Developer Cloud GPU pod.",
        "Automatic fallback keeps the demo stable if Gemma is unreachable.",
        "Containerized with Docker and built with Next.js + Supabase + Tailwind.",
        "Competes in Unicorn Track and Best AMD-Hosted Gemma Project."
    ])

# 9. Demo flow
add_content_slide(prs,
    "Live Demo Flow",
    [
        "1. Student types messy notes: 'photosynthesis chlorophyll sunlight CO2 oxygen glucose'.",
        "2. PADAYON detects Science → Biology → Photosynthesis and aligns it to Grade 9 curriculum.",
        "3. Clean notes, reviewer, flashcards, quiz, summary, and story are created automatically.",
        "4. Student asks in Cebuano and gets a Cebuano-first explanation with English terms.",
        "5. Quiz answer updates mastery score and learner profile.",
        "6. Student returns later: 'show my flashcards' — saved materials are retrieved instantly."
    ],
    image_path=os.path.join(IMG_DIR, "padayon_chat.png"),
    image_left=True)

# 10. Market potential
add_content_slide(prs,
    "Market Potential",
    [
        "Target: 27+ million Filipino K–12 students, especially those in Cebuano/Filipino-speaking regions.",
        "Go-to-market: public school pilots, review centers, and homeschool communities.",
        "Differentiation: not another answer bot — it builds study habits and long-term mastery.",
        "Expansion: add Araling Panlipunan, MAPEH, and senior high tracks; teacher dashboard for classroom analytics.",
        "Monetization: freemium student plan + school/institution licenses."
    ])

# 11. Why we can win
add_content_slide(prs,
    "Why PADAYON Can Win",
    [
        "Strong problem-solution fit backed by real PISA and language-barrier research.",
        "Complete, working demo — not a concept slide. Agents, memory, library, and Gemma toggle all functional.",
        "Clear use of AMD/Fireworks infrastructure with Gemma 4 integration.",
        "Differentiated in a sea of chatbots: organizes, teaches, remembers.",
        "Built for the Philippine context — translanguaging is a genuine competitive moat."
    ])

# 12. Links
add_content_slide(prs,
    "Try It Now",
    [
        "Live demo: https://courtesy-bacon-post-internet.trycloudflare.com",
        "GitHub: https://github.com/pjmorales1123/padayon",
        "Track: AMD Developer Hackathon: ACT II — Unicorn Track",
        "Also competing for: Best AMD-Hosted Gemma Project",
        "Contact: Prince Jireh Morales"
    ],
    subtitle="Built with Next.js, Supabase, Tailwind, Fireworks AI, and Gemma on AMD infrastructure")

prs.save(OUT)
print(f"Saved: {OUT}")
